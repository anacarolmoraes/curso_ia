import os
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

BRAND_BG = "ffffff"
BRAND_BG_ALT = "f9fafb"
BRAND_TEXT = "111827"
BRAND_TEXT_SECONDARY = "4b5563"
BRAND_ACCENT = "01467d"
BRAND_ACCENT_SECONDARY = "0078d4"
BRAND_ACCENT_TERTIARY = "ffc107"
BRAND_HEADING_FONT = "Montserrat"
BRAND_BODY_FONT = "Inter"

def apply_background(slide, bg_color=BRAND_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(bg_color)

def add_title_slide(prs, headline, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    top_bar.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(11.83), Inches(1.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = headline
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.6), Inches(10.33), Inches(1.0))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(24)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
        p.alignment = PP_ALIGN.CENTER

def add_quote_slide(prs, quote, attribution):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    
    quote_mark = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2.0), Inches(2.0))
    p = quote_mark.text_frame.paragraphs[0]
    p.text = "\u201C"
    p.font.name = "Georgia"
    p.font.size = Pt(200)
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT_TERTIARY)
    
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.8), Inches(0.12), Inches(2.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    accent_bar.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10.5), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(36)
    p.font.italic = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    
    attr = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10.5), Inches(0.6))
    p = attr.text_frame.paragraphs[0]
    p.text = attribution
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

def add_two_column_slide(prs, headline, left_title, left_points, right_title, right_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    p.text = headline
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    
    for idx, (title, points, left_inch) in enumerate([(left_title, left_points, 1.0), (right_title, right_points, 7.0)]):
        col_title = slide.shapes.add_textbox(Inches(left_inch), Inches(2.2), Inches(5.3), Inches(0.6))
        p = col_title.text_frame.paragraphs[0]
        p.text = title
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(28)
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
        
        col_body = slide.shapes.add_textbox(Inches(left_inch), Inches(3.0), Inches(5.3), Inches(3.5))
        tf = col_body.text_frame
        tf.word_wrap = True
        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.name = BRAND_BODY_FONT
            p.font.size = Pt(20)
            p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
            p.space_before = Pt(12)

def add_floating_cards_slide(prs, headline, cards):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, BRAND_BG_ALT)
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    p.text = headline
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    
    card_width = 3.5
    spacing = 0.4
    start_left = 1.0
    
    for idx, card in enumerate(cards):
        left = start_left + (idx * (card_width + spacing))
        
        # Card background
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.5), Inches(card_width), Inches(3.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)
        rect.line.fill.solid()
        rect.line.fill.fore_color.rgb = hex_to_rgb("e5e7eb")
        
        # Top accent
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(2.5), Inches(card_width), Inches(0.1))
        accent.fill.solid()
        accent.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_SECONDARY)
        accent.line.fill.background()
        
        # Card title
        tbox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(2.8), Inches(card_width - 0.4), Inches(0.6))
        p = tbox.text_frame.paragraphs[0]
        p.text = card["title"]
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(24)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
        
        # Card text
        body = slide.shapes.add_textbox(Inches(left + 0.2), Inches(3.5), Inches(card_width - 0.4), Inches(2.3))
        body.text_frame.word_wrap = True
        p = body.text_frame.paragraphs[0]
        p.text = card["text"]
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

def add_section_break(prs, text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, BRAND_ACCENT)
    
    box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb("ffffff")
    p.alignment = PP_ALIGN.CENTER
    
def add_stats_slide(prs, headline, stats):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    p.text = headline
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    
    stat_width = 4.0
    start_left = 2.0
    spacing = 1.3
    
    for idx, stat in enumerate(stats):
        left = start_left + (idx * (stat_width + spacing))
        
        # Number
        num = slide.shapes.add_textbox(Inches(left), Inches(3.0), Inches(stat_width), Inches(1.5))
        p = num.text_frame.paragraphs[0]
        p.text = stat["number"]
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(80)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
        p.alignment = PP_ALIGN.CENTER
        
        # Text
        desc = slide.shapes.add_textbox(Inches(left), Inches(4.5), Inches(stat_width), Inches(1.0))
        p = desc.text_frame.paragraphs[0]
        p.text = stat["text"]
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
        p.alignment = PP_ALIGN.CENTER

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 1. Title
    add_title_slide(prs, "Superávit de Produtividade", "A Inteligência Artificial Aplicada à Execução Orçamentária")
    
    # 2. Quote - Pain
    add_quote_slide(prs, "Fazer isso na mão é o nosso vilão. Rouba nosso tempo técnico para fazer trabalho braçal.", "— O Ponto de Dor do Contador")
    
    # 3. Two column - Paradigm shift
    add_two_column_slide(prs, "O Novo Paradigma", 
        "O Modelo Antigo", ["Busca manual", "PDFs de 8MB inteiros", "Sobrecarga de contexto", "Resultados genéricos e erros"],
        "A Delegação Tecnológica", ["Operação de Agente Autônomo", "Limpeza e curadoria", "Economia de Tokens", "Precisão e aderência normativa"])
        
    # 4. Floating Cards - Orquestração
    add_floating_cards_slide(prs, "O Método de Orquestração", [
        {"title": "Entrada", "text": "Curadoria e Limpeza de Dados (PDF para Markdown)."},
        {"title": "Processamento", "text": "Aplicação da Lógica Contábil e Regulatória no Prompt."},
        {"title": "Saída", "text": "Relatório preciso com validação da Classificação Orçamentária."}
    ])
    
    # 5. Section Break
    add_section_break(prs, "Hora 2: A Visão de Futuro")
    
    # 6. Stats - O Impacto
    add_stats_slide(prs, "O Impacto na Prática", [
        {"number": "2 Hrs", "text": "De trabalho braçal"},
        {"number": "40 Seg", "text": "Para auditoria com IA"}
    ])
    
    # 7. Quote
    add_quote_slide(prs, "A tela preta é como a cozinha de um restaurante no horário de pico. Você foca no 'Pedido' e a IA entrega o 'Prato'.", "— A Vacina Contra o Pânico")
    
    # 8. Two column - Delegation
    add_two_column_slide(prs, "Saneamento de Dados",
        "PDF Bruto", ["8MB de tamanho", "Lixo visual e logomarcas", "Causa alucinações na IA", "Desperdício de Tokens"],
        "Texto Limpo (Markdown)", ["Arquivo leve", "Apenas a informação central", "Foco no filé mignon do dado", "Economia do orçamento de atenção"])
        
    # 9. Section Break
    add_section_break(prs, "Hora 3: A Realidade (Segunda-feira)")
    
    # 10. Floating Cards - Survival Guide
    add_floating_cards_slide(prs, "O Guia de Sobrevivência", [
        {"title": "O Kit do Aluno", "text": "Um bloco de notas com o Prompt Perfeito. Sua nova Skill portátil."},
        {"title": "Curadoria Manual", "text": "Se não há ferramenta de conversão, copie apenas o texto necessário."},
        {"title": "Navegador Web", "text": "Usando o Google Gemini ou Claude.ai para executar a tarefa."}
    ])
    
    # 11. Quote
    add_quote_slide(prs, "A IA não substitui a assinatura ou a responsabilidade. Ela substitui o trabalho braçal.", "— O Futuro da Auditoria")
    
    # 12. Title (Closing)
    add_title_slide(prs, "Vamos Orquestrar?", "Obrigado!")
    
    output_dir = Path("output/tceto")
    output_dir.mkdir(parents=True, exist_ok=True)
    prs.save(output_dir / "tceto-curso-ia.pptx")
    print(f"Presentation saved to {output_dir / 'tceto-curso-ia.pptx'}")

if __name__ == "__main__":
    main()
